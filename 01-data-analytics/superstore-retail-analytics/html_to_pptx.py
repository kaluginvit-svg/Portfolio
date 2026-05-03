# -*- coding: utf-8 -*-
"""
Конвертация HTML-презентации в PPTX: каждая страница рендерится в изображение
и вставляется в слайд в полный размер. Требуется: pip install playwright python-pptx
и один раз: playwright install chromium
"""
import os
import sys
import tempfile
import time

# Размер viewport при съёмке (16:9) — должен совпадать с размером слайда в PPTX
VIEWPORT_W = 1280
VIEWPORT_H = 720

# Пауза после переключения слайда (сек): дождаться завершения анимации
# (переход слайда 0.7s + reveal-элементы с задержкой + bar-fill 1.5s)
SLIDE_ANIMATION_WAIT = 3.0


def capture_slides(html_path, output_dir, total_slides=10):
    """Рендерит каждый слайд HTML в PNG через Playwright."""
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        print("Установите Playwright: pip install playwright")
        print("Затем: playwright install chromium")
        sys.exit(1)

    html_abs = os.path.abspath(html_path)
    if not os.path.isfile(html_abs):
        print("Файл не найден:", html_abs)
        sys.exit(1)

    file_url = "file:///" + html_abs.replace("\\", "/")

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page()
        page.set_viewport_size({"width": VIEWPORT_W, "height": VIEWPORT_H})
        page.goto(file_url, wait_until="networkidle", timeout=30000)
        # Дождаться применения стилей и шрифтов
        time.sleep(1)
        # Скрыть панель навигации, чтобы в кадр попадал только слайд
        page.evaluate("() => { const el = document.querySelector('.nav-controls'); if (el) el.style.display = 'none'; }")

        for i in range(1, total_slides + 1):
            page.evaluate(f"goToSlide({i})")
            time.sleep(SLIDE_ANIMATION_WAIT)  # дождаться завершения анимации слайда
            path = os.path.join(output_dir, f"slide_{i:02d}.png")
            page.screenshot(path=path, type="png")
            print("  Слайд", i, "->", path)

        browser.close()


# Слайд с ссылкой на дашборд (1-based) и имя файла по умолчанию.
# Чтобы ссылка работала, superstore-dashboard.html должен лежать в той же папке, что и .pptx.
DASHBOARD_SLIDE_INDEX = 4  # 4-й слайд в презентации
DASHBOARD_HTML_NAME = "superstore-dashboard.html"


def _file_uri(path):
    """Абсолютный file:// URL для пути (чтобы ссылка в PowerPoint открывала файл)."""
    abs_path = os.path.abspath(path)
    # file:///C:/path/form — PowerPoint надёжно открывает по такому формату
    return "file:///" + abs_path.replace("\\", "/")


def _disable_advance_on_click(slide):
    """Отключает «переход по клику» у слайда (OOXML: p:transition advClick="0")."""
    from lxml import etree
    NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    el = slide._element
    trans = el.find(f".//{{{NS}}}transition")
    if trans is None:
        trans = etree.SubElement(el, f"{{{NS}}}transition")
    trans.set("advClick", "0")


def build_pptx(image_paths, output_pptx, dashboard_link_path=None):
    """Создаёт PPTX: один слайд на изображение, картинка на весь слайд.
    На слайде DASHBOARD_SLIDE_INDEX — кнопка с гиперссылкой; у этого слайда отключён переход по клику.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    dashboard_slide = None  # слайд с кнопкой дашборда

    for idx, path in enumerate(image_paths):
        if not os.path.isfile(path):
            continue
        slide = prs.slides.add_slide(blank)
        pic = slide.shapes.add_picture(
            path,
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
        if (idx + 1) == DASHBOARD_SLIDE_INDEX and dashboard_link_path:
            uri = dashboard_link_path if dashboard_link_path.startswith("file:") else _file_uri(dashboard_link_path)
            btn = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(5.0),
                Inches(5.9),
                Inches(3.3),
                Inches(0.65),
            )
            btn.fill.solid()
            btn.fill.fore_color.rgb = RGBColor(59, 130, 246)
            btn.line.color.rgb = RGBColor(255, 255, 255)
            btn.line.width = Pt(0.5)
            # Гиперссылка: используем .hyperlink.address (не .hyperlink_address)
            btn.click_action.hyperlink.address = uri
            tf = btn.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Открыть дашборд"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            dashboard_slide = slide

    # Отключаем переход по клику на слайде с дашбордом — тогда клик по кнопке откроет ссылку
    if dashboard_slide is not None:
        _disable_advance_on_click(dashboard_slide)

    prs.save(output_pptx)
    print("Сохранено:", output_pptx)


def main():
    html_name = "superstore_presentation.html"
    pptx_name = "superstore_presentation.pptx"
    script_dir = os.path.dirname(os.path.abspath(__file__))
    html_path = os.path.join(script_dir, html_name)
    pptx_path = os.path.join(script_dir, pptx_name)

    # Количество слайдов в HTML (по числу .slide)
    total_slides = 10

    with tempfile.TemporaryDirectory(prefix="html2pptx_") as tmpdir:
        print("Рендер слайдов в изображения...")
        capture_slides(html_path, tmpdir, total_slides)

        image_paths = [
            os.path.join(tmpdir, f"slide_{i:02d}.png")
            for i in range(1, total_slides + 1)
        ]
        # Ссылка на дашборд: абсолютный путь → в PPTX попадёт file:// URL
        dashboard_path = os.path.join(script_dir, DASHBOARD_HTML_NAME)
        print("Сборка PPTX...")
        build_pptx(image_paths, pptx_path, dashboard_link_path=dashboard_path)

    print("Готово.")


if __name__ == "__main__":
    main()
