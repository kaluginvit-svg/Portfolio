from __future__ import annotations

from collections import Counter
from pathlib import Path

from pptx import Presentation


def rgb_to_hex(rgb) -> str | None:
    if rgb is None:
        return None
    try:
        return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    except Exception:
        return None


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    pptx_path = root / "financial_practice_brandbook.pptx"
    prs = Presentation(str(pptx_path))

    fills = Counter()
    lines = Counter()
    texts = Counter()

    for slide in prs.slides:
        for shape in slide.shapes:
            # fill
            try:
                fill = shape.fill
                if fill and fill.type is not None and getattr(fill, "fore_color", None):
                    fc = fill.fore_color
                    hx = rgb_to_hex(getattr(fc, "rgb", None))
                    if hx:
                        fills[hx] += 1
            except Exception:
                pass

            # line
            try:
                ln = shape.line
                if ln and getattr(ln, "color", None):
                    hx = rgb_to_hex(getattr(ln.color, "rgb", None))
                    if hx:
                        lines[hx] += 1
            except Exception:
                pass

            # text
            if getattr(shape, "has_text_frame", False):
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if not (run.text or "").strip():
                            continue
                        try:
                            c = run.font.color
                            hx = rgb_to_hex(getattr(getattr(c, "rgb", None), "__iter__", None) and c.rgb)  # type: ignore
                        except Exception:
                            hx = None
                        if hx:
                            texts[hx] += 1

    print("top fills:", fills.most_common(12))
    print("top lines:", lines.most_common(12))
    print("top text :", texts.most_common(12))


if __name__ == "__main__":
    main()

