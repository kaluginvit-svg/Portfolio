from __future__ import annotations

from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation


@dataclass(frozen=True)
class SlideStats:
    idx: int
    shape_count: int
    text_shape_count: int
    font_sizes: Counter[int]
    sample_texts: list[str]


def iter_shape_texts(slide) -> Iterable[tuple[str, int | None]]:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        for p in tf.paragraphs:
            for run in p.runs:
                t = (run.text or "").strip()
                if not t:
                    continue
                size = int(run.font.size.pt) if run.font.size else None
                yield t, size


def collect_slide_stats(prs: Presentation) -> list[SlideStats]:
    out: list[SlideStats] = []
    for i, slide in enumerate(prs.slides, start=1):
        sizes: Counter[int] = Counter()
        sample: list[str] = []
        text_shape_count = 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text_shape_count += 1
        for t, size in iter_shape_texts(slide):
            if size is not None:
                sizes[size] += 1
            if len(sample) < 10:
                sample.append(t)
        out.append(
            SlideStats(
                idx=i,
                shape_count=len(slide.shapes),
                text_shape_count=text_shape_count,
                font_sizes=sizes,
                sample_texts=sample,
            )
        )
    return out


def main() -> None:
    here = Path(__file__).resolve()
    root = here.parents[1]
    pptx_path = root / "financial_practice_brandbook.pptx"
    prs = Presentation(str(pptx_path))

    print(f"file: {pptx_path}")
    print(f"slides: {len(prs.slides)}")
    print()

    for st in collect_slide_stats(prs):
        common = ", ".join(f"{k}ptx{v}" for k, v in st.font_sizes.most_common(8))
        sample = " | ".join(st.sample_texts[:6])
        print(
            f"{st.idx:02d}: shapes={st.shape_count}, text_shapes={st.text_shape_count}, "
            f"fonts=[{common}]"
        )
        if sample:
            print(f"    sample: {sample}")


if __name__ == "__main__":
    main()

