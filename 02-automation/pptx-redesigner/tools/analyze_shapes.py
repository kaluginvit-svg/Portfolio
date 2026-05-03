from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    pptx_path = root / "financial_practice_brandbook.pptx"
    prs = Presentation(str(pptx_path))

    print(f"file: {pptx_path}")
    print(f"slides: {len(prs.slides)}")
    print()

    for si, slide in enumerate(prs.slides, start=1):
        counts = {}
        for shape in slide.shapes:
            counts[shape.shape_type] = counts.get(shape.shape_type, 0) + 1
        def fmt(st):
            name = str(st)
            return name.split(".")[-1] if "." in name else name
        top = ", ".join(f"{fmt(k)}={v}" for k, v in sorted(counts.items(), key=lambda kv: -kv[1])[:8])
        print(f"{si:02d}: {top}")

        # Heuristic: list small non-text shapes that could be icons
        iconish = []
        for idx, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_text_frame", False):
                continue
            if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.GROUP):
                continue
            w = shape.width
            h = shape.height
            if w is None or h is None:
                continue
            # under ~1.2 inch square-ish
            if w < 1100000 and h < 1100000:
                iconish.append((idx, shape.shape_type, int(w), int(h)))
        if iconish:
            print("    iconish:", ", ".join(f"#{i}:{str(t).split('.')[-1]} {w}x{h}" for i,t,w,h in iconish[:12]))


if __name__ == "__main__":
    main()

