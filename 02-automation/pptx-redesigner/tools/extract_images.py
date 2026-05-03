from __future__ import annotations

import hashlib
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def sha1(data: bytes) -> str:
    return hashlib.sha1(data).hexdigest()


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    pptx_path = root / "financial_practice_brandbook.pptx"
    out_dir = root / "extracted_assets"
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    seen: set[str] = set()
    total = 0

    for si, slide in enumerate(prs.slides, start=1):
        for shi, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            img = shape.image
            blob = img.blob
            h = sha1(blob)
            if h in seen:
                continue
            seen.add(h)
            total += 1
            ext = img.ext or "bin"
            name = f"s{si:02d}_p{shi:02d}_{h[:10]}.{ext}"
            (out_dir / name).write_bytes(blob)

    print(f"exported_unique_images: {total}")
    print(f"dir: {out_dir}")


if __name__ == "__main__":
    main()

