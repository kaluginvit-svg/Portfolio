from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def emu_to_in(x: int) -> float:
    return float(Emu(x).inches)


def shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    out = []
    for p in shape.text_frame.paragraphs:
        t = "".join(r.text for r in p.runs).strip()
        if t:
            out.append(t)
    return " / ".join(out)


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    pptx_path = root / "financial_practice_brandbook_redesign.pptx"
    prs = Presentation(str(pptx_path))

    # Make stdout UTF-8 on Windows consoles that default to legacy encodings
    try:
        import sys

        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass

    print(f"file: {pptx_path}")
    print(f"slides: {len(prs.slides)}")
    print(f"size: {emu_to_in(prs.slide_width):.2f}in x {emu_to_in(prs.slide_height):.2f}in")
    print()

    for si, slide in enumerate(prs.slides, start=1):
        print(f"--- slide {si:02d} ---")
        rows = []
        for idx, sh in enumerate(slide.shapes, start=1):
            try:
                l, t, w, h = sh.left, sh.top, sh.width, sh.height
            except Exception:
                continue
            txt = shape_text(sh)
            if txt:
                # top font size sample
                fs = []
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if (r.text or "").strip() and r.font.size:
                            fs.append(int(r.font.size.pt))
                fs_s = ",".join(map(str, sorted(set(fs))[:6]))
            else:
                fs_s = ""
            rows.append(
                (
                    t,
                    idx,
                    f"{emu_to_in(l):.2f},{emu_to_in(t):.2f} {emu_to_in(w):.2f}x{emu_to_in(h):.2f}",
                    ("T " + fs_s) if txt else "—",
                    (txt[:90] + ("…" if len(txt) > 90 else "")) if txt else "",
                )
            )
        rows.sort(key=lambda r: (r[0], r[2]))
        for _, idx, box, kind, txt in rows[:60]:
            print(f"{idx:02d} {box:>20} {kind:<10} {txt}")
        if len(rows) > 60:
            print(f"... ({len(rows)-60} more shapes)")
        print()


if __name__ == "__main__":
    main()

