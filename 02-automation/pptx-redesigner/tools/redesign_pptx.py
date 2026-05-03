from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


BRAND_BURGUNDY = RGBColor(0x56, 0x00, 0x04)  # #560004
BRAND_NAVY = RGBColor(0x04, 0x04, 0x42)  # #040442
BRAND_MUTED = RGBColor(0x84, 0x38, 0x3E)  # #84383E
BRAND_TEXT = RGBColor(0x0F, 0x0F, 0x0F)  # #0F0F0F
BRAND_TEXT_MUTED = RGBColor(0x54, 0x54, 0x54)  # #545454
BRAND_LINE = RGBColor(0xD6, 0xD6, 0xD6)  # #D6D6D6


@dataclass(frozen=True)
class FontRules:
    title_min: int = 34
    subtitle_min: int = 22
    body_min: int = 18
    small_min: int = 12


def _set_run_font(run, pt: int, *, bold: bool | None = None, color: RGBColor | None = None) -> None:
    run.font.size = Pt(pt)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def upscale_text(slide, rules: FontRules) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                txt = (run.text or "").strip()
                if not txt:
                    continue

                # heuristic: slide index markers like "02" stay small
                if txt.isdigit() and 1 <= len(txt) <= 2:
                    _set_run_font(run, max(rules.small_min, 12), bold=True, color=BRAND_MUTED)
                    continue

                current = int(run.font.size.pt) if run.font.size else None
                if current is None:
                    # default body
                    _set_run_font(run, rules.body_min, color=BRAND_TEXT)
                    continue

                # map existing sizes to new hierarchy
                if current >= 26:
                    _set_run_font(run, max(rules.title_min, current + 8), bold=True, color=BRAND_TEXT)
                elif current >= 22:
                    _set_run_font(run, max(rules.subtitle_min, current + 4), bold=True, color=BRAND_TEXT)
                elif current >= 18:
                    _set_run_font(run, max(rules.body_min, current + 2), color=BRAND_TEXT)
                elif current >= 14:
                    _set_run_font(run, max(rules.body_min, 18), color=BRAND_TEXT_MUTED)
                else:
                    _set_run_font(run, max(rules.body_min, 18), color=BRAND_TEXT_MUTED)


def tighten_vertical(slide, slide_w: int, slide_h: int, *, top_pad_in: float = 0.45, bottom_pad_in: float = 0.35) -> None:
    # Shift content upward if there's too much top padding.
    # Skip background-like big shapes that cover most slide.

    tops = []
    bottoms = []
    for shape in slide.shapes:
        try:
            if shape.width >= int(slide_w * 0.95) and shape.height >= int(slide_h * 0.90):
                continue
        except Exception:
            pass
        try:
            tops.append(shape.top)
            bottoms.append(shape.top + shape.height)
        except Exception:
            continue

    if not tops:
        return

    min_top = min(tops)
    max_bottom = max(bottoms)
    desired_top = Inches(top_pad_in)
    desired_bottom = slide_h - Inches(bottom_pad_in)

    dy_up = min_top - desired_top
    dy_down = desired_bottom - max_bottom

    # If content starts too low -> move up by part of delta (not full, to avoid clipping)
    if dy_up > Inches(0.15):
        shift = min(dy_up, Inches(0.35))
        for shape in slide.shapes:
            try:
                shape.top = max(0, shape.top - shift)
            except Exception:
                pass

    # If bottom is far from edge -> expand a bit downward by nudging lower elements
    if dy_down > Inches(0.35):
        # nudge shapes that are in bottom half
        for shape in slide.shapes:
            try:
                if shape.top > slide_h * 0.45:
                    shape.top = min(slide_h - shape.height, shape.top + Inches(0.15))
            except Exception:
                pass


def _style_icon_line(shape, *, color: RGBColor) -> None:
    shape.fill.background()  # no fill
    shape.line.color.rgb = color
    shape.line.width = Pt(1.75)


def add_icon_growth(slide, left, top, size, *, color: RGBColor = BRAND_BURGUNDY) -> None:
    # simple "chart up": baseline + 3 bars + arrow
    # Use thin rectangles instead of LINE (python-pptx doesn't expose a LINE autoshape).
    base = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top + int(size * 0.85), int(size), int(size * 0.04)
    )
    base.fill.solid()
    base.fill.fore_color.rgb = color
    base.line.fill.background()

    for i, hmul in enumerate([0.35, 0.55, 0.75]):
        bar_w = int(size * 0.16)
        gap = int(size * 0.08)
        x = left + i * (bar_w + gap)
        y = top + int(size * (1 - hmul))
        rect = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, bar_w, int(size * hmul)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()

    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, left + int(size * 0.55), top + int(size * 0.10), int(size * 0.42), int(size * 0.28)
    )
    arrow.fill.background()
    arrow.line.color.rgb = color
    arrow.line.width = Pt(1.75)


def add_icon_control(slide, left, top, size, *, color: RGBColor = BRAND_BURGUNDY) -> None:
    # shield outline + check
    shield = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEART, left, top, int(size * 0.95), int(size * 0.95))
    shield.fill.background()
    shield.line.color.rgb = color
    shield.line.width = Pt(1.75)
    # check: two thin rectangles
    l1 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left + int(size * 0.30),
        top + int(size * 0.62),
        int(size * 0.22),
        int(size * 0.06),
    )
    l1.rotation = 45
    l1.fill.solid()
    l1.fill.fore_color.rgb = color
    l1.line.fill.background()
    l2 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left + int(size * 0.42),
        top + int(size * 0.58),
        int(size * 0.34),
        int(size * 0.06),
    )
    l2.rotation = -45
    l2.fill.solid()
    l2.fill.fore_color.rgb = color
    l2.line.fill.background()


def add_icon_money(slide, left, top, size, *, color: RGBColor = BRAND_BURGUNDY) -> None:
    coin = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, size, size)
    coin.fill.background()
    coin.line.color.rgb = color
    coin.line.width = Pt(1.75)
    s = slide.shapes.add_textbox(left, top + int(size * 0.12), size, int(size * 0.76))
    tf = s.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "₽"
    r.font.bold = True
    r.font.size = Pt(20)
    r.font.color.rgb = color
    p.alignment = 1  # center


def add_icon_analytics(slide, left, top, size, *, color: RGBColor = BRAND_BURGUNDY) -> None:
    # magnifier: circle + handle
    c = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, int(size * 0.72), int(size * 0.72))
    c.fill.background()
    c.line.color.rgb = color
    c.line.width = Pt(1.75)
    h = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left + int(size * 0.54),
        top + int(size * 0.58),
        int(size * 0.34),
        int(size * 0.06),
    )
    h.rotation = 45
    h.fill.solid()
    h.fill.fore_color.rgb = color
    h.line.fill.background()


def add_icons_per_slide(slide, idx: int) -> None:
    # Place small icons near key blocks without relying on exact layout.
    # Heuristic anchors: top-left region under title.
    size = int(Inches(0.35))
    x0 = int(Inches(0.65))
    y0 = int(Inches(1.55))
    dx = int(Inches(0.55))

    if idx == 2:
        add_icon_growth(slide, x0, y0, size, color=BRAND_BURGUNDY)
        add_icon_control(slide, x0 + dx, y0, size, color=BRAND_BURGUNDY)
        add_icon_money(slide, x0 + dx * 2, y0, size, color=BRAND_BURGUNDY)
        add_icon_analytics(slide, x0 + dx * 3, y0, size, color=BRAND_BURGUNDY)
    elif idx == 4:
        add_icon_analytics(slide, x0, y0, size, color=BRAND_BURGUNDY)
        add_icon_control(slide, x0 + dx, y0, size, color=BRAND_BURGUNDY)
        add_icon_money(slide, x0 + dx * 2, y0, size, color=BRAND_BURGUNDY)
    elif idx in (6, 7):
        add_icon_analytics(slide, x0, y0, size, color=BRAND_MUTED)
        add_icon_control(slide, x0 + dx, y0, size, color=BRAND_MUTED)
    elif idx == 8:
        add_icon_growth(slide, x0, y0, size, color=BRAND_NAVY)
        add_icon_money(slide, x0 + dx, y0, size, color=BRAND_NAVY)
        add_icon_control(slide, x0 + dx * 2, y0, size, color=BRAND_NAVY)
    elif idx in (10, 11):
        add_icon_analytics(slide, x0, y0, size, color=BRAND_MUTED)


def style_cards(slide) -> None:
    # Make rectangles feel tighter: white fill, subtle border, no shadows (avoid effects).
    for shape in slide.shapes:
        try:
            if shape.shape_type != 1:  # AUTO_SHAPE
                continue
            if not hasattr(shape, "auto_shape_type"):
                continue
            if shape.auto_shape_type != MSO_AUTO_SHAPE_TYPE.RECTANGLE:
                continue
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            shape.line.color.rgb = BRAND_LINE
            shape.line.width = Pt(1.0)
        except Exception:
            continue


def add_subtle_watermark_on_title(slide, prs: Presentation) -> None:
    # Slide 1: add a very light brand mark as watermark (use extracted logo mark image).
    root = Path(__file__).resolve().parents[1]
    mark_path = root / "extracted_assets" / "s01_p02_3c38d998e3.png"
    if not mark_path.exists():
        return
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    pic = slide.shapes.add_picture(str(mark_path), int(slide_w * 0.60), int(slide_h * 0.10), width=int(slide_w * 0.35))
    # transparency isn't directly supported; mimic by putting it behind and small.
    try:
        pic.crop_left = 0
    except Exception:
        pass


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    src = root / "financial_practice_brandbook.pptx"
    dst = root / "financial_practice_brandbook_redesign.pptx"

    prs = Presentation(str(src))
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    rules = FontRules()
    for idx, slide in enumerate(prs.slides, start=1):
        upscale_text(slide, rules)
        style_cards(slide)
        tighten_vertical(slide, slide_w, slide_h)
        add_icons_per_slide(slide, idx)
        if idx == 1:
            add_subtle_watermark_on_title(slide, prs)

    prs.save(str(dst))
    print(f"saved: {dst}")


if __name__ == "__main__":
    main()

