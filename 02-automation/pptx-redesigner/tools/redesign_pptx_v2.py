from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# Brand palette (from your spec)
BRAND_BURGUNDY = RGBColor(0x56, 0x00, 0x04)  # #560004
BRAND_DARK = RGBColor(0x0F, 0x0F, 0x0F)  # #0F0F0F
BRAND_GRAY = RGBColor(0x54, 0x54, 0x54)  # #545454
BRAND_LIGHT_GRAY = RGBColor(0xD6, 0xD6, 0xD6)  # #D6D6D6
BRAND_NAVY = RGBColor(0x04, 0x04, 0x42)  # #040442
BRAND_WINE = RGBColor(0x84, 0x38, 0x3E)  # #84383E


@dataclass(frozen=True)
class Scale:
    title: float
    subtitle: float
    card_title: float
    body: float
    bullets: float
    footer: float


SLIDE1 = Scale(title=28, subtitle=13.5, card_title=18, body=12.5, bullets=14, footer=10.5)
DEFAULT = Scale(title=26, subtitle=11, card_title=18, body=12, bullets=13.5, footer=10.5)


def _iter_runs(shape):
    if not getattr(shape, "has_text_frame", False):
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            yield p, r


def _set_shape_margins(shape, *, left=0.12, right=0.10, top=0.06, bottom=0.06) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.margin_left = Inches(left)
    tf.margin_right = Inches(right)
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(bottom)


def _tighten_paragraph(p) -> None:
    # Remove extra spacing to make composition denser
    try:
        p.space_before = Pt(0)
        p.space_after = Pt(0)
    except Exception:
        pass


def _apply_font(shape, pt: float, *, bold: bool | None = None, color: RGBColor | None = None) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    for p, r in _iter_runs(shape):
        if not (r.text or "").strip():
            continue
        r.font.size = Pt(pt)
        if bold is not None:
            r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color
        _tighten_paragraph(p)


def _shape_max_font_pt(shape) -> int | None:
    mx = None
    if not getattr(shape, "has_text_frame", False):
        return None
    for _, r in _iter_runs(shape):
        if (r.text or "").strip() and r.font.size:
            v = int(r.font.size.pt)
            mx = v if mx is None else max(mx, v)
    return mx


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    parts = []
    for p in shape.text_frame.paragraphs:
        t = "".join(r.text for r in p.runs).strip()
        if t:
            parts.append(t)
    return "\n".join(parts).strip()


def classify_and_restylize_slide(slide, slide_idx: int, scale: Scale) -> None:
    # Identify title/subtitle by position and current size (since we must preserve structure).
    title_candidates = []
    subtitle_candidates = []
    footer_candidates = []

    all_text_shapes = []
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        txt = _shape_text(sh)
        if not txt:
            continue
        mx = _shape_max_font_pt(sh) or 0
        all_text_shapes.append((mx, sh))
        top_in = float(sh.top) / 914400  # EMU per inch
        h_in = float(sh.height) / 914400

        # Page number marker like "02"
        if txt.strip().isdigit() and len(txt.strip()) <= 2 and mx <= 14:
            _apply_font(sh, scale.footer, bold=True, color=BRAND_WINE)
            continue

        # title: very top, widest, biggest font
        if top_in <= 0.55 and mx >= 26 and h_in <= 0.7:
            title_candidates.append((mx, -sh.width, sh))
            continue

        # subtitle: under title row, medium size
        if top_in <= 1.05 and 14 <= mx <= 22 and h_in <= 0.35:
            subtitle_candidates.append((mx, -sh.width, sh))
            continue

        # footer/service line: bottom thin line
        if top_in >= 6.3 and h_in <= 0.35 and mx >= 12:
            footer_candidates.append((mx, sh))

    # Fallback: on some layouts (esp. slide 1) the title isn't in the very top band.
    if not title_candidates and all_text_shapes:
        mx, sh = sorted(all_text_shapes, key=lambda x: x[0], reverse=True)[0]
        if mx >= 26:
            title_candidates.append((mx, -sh.width, sh))

    if title_candidates:
        _, __, sh = sorted(title_candidates, reverse=True)[0]
        _apply_font(sh, scale.title, bold=True, color=BRAND_DARK)
        # make title not too airy
        _set_shape_margins(sh, top=0.04, bottom=0.02)

    if subtitle_candidates:
        _, __, sh = sorted(subtitle_candidates, reverse=True)[0]
        _apply_font(sh, scale.subtitle, bold=False, color=BRAND_GRAY)
        _set_shape_margins(sh, top=0.02, bottom=0.02)

    for _, sh in footer_candidates:
        _apply_font(sh, scale.footer, bold=False, color=BRAND_GRAY)

    # Restyle the rest: bring card headings and body to target sizes.
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        txt = _shape_text(sh)
        if not txt:
            continue

        mx = _shape_max_font_pt(sh) or 0
        top_in = float(sh.top) / 914400
        h_in = float(sh.height) / 914400

        # skip title/subtitle already handled
        if top_in <= 1.05 and mx >= 22:
            continue

        # Card headers: short height, stronger size
        if h_in <= 0.45 and mx >= 16:
            _apply_font(sh, scale.card_title, bold=True, color=BRAND_DARK)
            _set_shape_margins(sh, top=0.03, bottom=0.00)
            continue

        # Lists: multiple lines or slash-separated patterns
        is_listish = ("\n" in txt) or (" / " in txt) or ("•" in txt) or ("-" in txt[:3])
        if is_listish:
            _apply_font(sh, scale.bullets, bold=False, color=BRAND_GRAY)
            _set_shape_margins(sh, top=0.06, bottom=0.06)
        else:
            _apply_font(sh, scale.body, bold=False, color=BRAND_GRAY)
            _set_shape_margins(sh, top=0.06, bottom=0.06)


def tighten_slide_bbox(prs: Presentation, slide) -> None:
    # Reduce excessive outer padding by fitting all content into a tighter safe area.
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    left_margin = Inches(0.55)
    right_margin = Inches(0.55)
    top_margin = Inches(0.40)
    bottom_margin = Inches(0.35)

    min_l, min_t = None, None
    max_r, max_b = None, None

    for sh in slide.shapes:
        try:
            # ignore full-width tiny top strip and similar decorative thin lines
            if sh.height <= Inches(0.22) and sh.width >= int(slide_w * 0.95):
                continue
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
        except Exception:
            continue
        min_l = l if min_l is None else min(min_l, l)
        min_t = t if min_t is None else min(min_t, t)
        max_r = (l + w) if max_r is None else max(max_r, l + w)
        max_b = (t + h) if max_b is None else max(max_b, t + h)

    if min_l is None:
        return

    src_w = max_r - min_l
    src_h = max_b - min_t
    dst_l = left_margin
    dst_t = top_margin
    dst_w = slide_w - left_margin - right_margin
    dst_h = slide_h - top_margin - bottom_margin

    # Scale slightly only if there is big unused padding; avoid enlarging beyond 1.08
    sx = float(dst_w) / float(src_w) if src_w else 1.0
    sy = float(dst_h) / float(src_h) if src_h else 1.0
    s = min(1.08, max(0.92, min(sx, sy)))

    # Translate to target margins after scaling around top-left bbox corner
    for sh in slide.shapes:
        try:
            if sh.height <= Inches(0.22) and sh.width >= int(slide_w * 0.95):
                continue
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
            nl = int(dst_l + (l - min_l) * s)
            nt = int(dst_t + (t - min_t) * s)
            nw = int(w * s)
            nh = int(h * s)
            sh.left, sh.top, sh.width, sh.height = nl, nt, nw, nh
        except Exception:
            continue


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    src = root / "financial_practice_brandbook_redesign.pptx"
    dst = root / "financial_practice_brandbook_redesign_v2.pptx"

    prs = Presentation(str(src))

    for i, slide in enumerate(prs.slides, start=1):
        scale = SLIDE1 if i == 1 else DEFAULT
        classify_and_restylize_slide(slide, i, scale)
        tighten_slide_bbox(prs, slide)

    prs.save(str(dst))
    print(f"saved: {dst}")


if __name__ == "__main__":
    main()

